import os
import re
import io
import json
import zipfile
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple

import streamlit as st
import pandas as pd
import numpy as np
from PIL import Image


# ============================================================
# CONFIG
# ============================================================
st.set_page_config(page_title="DocFlow Converter", layout="wide")

# ============================================================
# PREMIUM UI (Website feel)
# ============================================================
st.markdown(
    """
    <style>
      /* Base */
      html, body, [class*="stApp"] {
        background: radial-gradient(1200px 800px at 10% 0%, rgba(99,102,241,0.18), transparent 55%),
                    radial-gradient(900px 700px at 85% 10%, rgba(16,185,129,0.14), transparent 55%),
                    linear-gradient(180deg, #f7f8fb, #f3f4f6);
        color: #0f172a;
      }
      .block-container { max-width: 1180px; padding-top: 1.2rem; padding-bottom: 2.0rem; }

      /* Hide Streamlit footer/menu */
      #MainMenu, footer { visibility: hidden; }
      header { visibility: hidden; }

      /* Hero */
      .hero {
        display:flex; align-items:center; justify-content:space-between;
        gap: 16px;
        padding: 16px 18px;
        border-radius: 18px;
        border: 1px solid rgba(15,23,42,0.08);
        background: rgba(255,255,255,0.75);
        backdrop-filter: blur(10px);
        box-shadow: 0 10px 24px rgba(15,23,42,0.06);
        margin-bottom: 14px;
      }
      .brandwrap { display:flex; align-items:center; gap: 12px; }
      .logo {
        width: 42px; height: 42px;
        border-radius: 14px;
        background: linear-gradient(135deg, rgba(99,102,241,0.95), rgba(16,185,129,0.90));
        box-shadow: 0 10px 18px rgba(99,102,241,0.18);
      }
      .brand {
        font-size: 18px; font-weight: 900; letter-spacing: 0.2px;
        margin: 0;
      }
      .tagline {
        margin-top: 2px;
        color: rgba(15,23,42,0.65);
        font-size: 13px;
      }
      .chiprow { display:flex; gap: 8px; flex-wrap: wrap; justify-content:flex-end; }
      .chip {
        font-size: 12px; padding: 6px 10px;
        border-radius: 999px;
        border: 1px solid rgba(15,23,42,0.10);
        background: rgba(255,255,255,0.70);
        color: rgba(15,23,42,0.75);
      }

      /* Cards */
      .card {
        border-radius: 18px;
        border: 1px solid rgba(15,23,42,0.08);
        background: rgba(255,255,255,0.72);
        backdrop-filter: blur(10px);
        box-shadow: 0 12px 28px rgba(15,23,42,0.06);
        padding: 14px 14px;
      }
      .cardtitle {
        display:flex; align-items:center; justify-content:space-between;
        margin-bottom: 8px;
      }
      .cardtitle h3 {
        margin: 0;
        font-size: 14px;
        font-weight: 900;
        color: rgba(15,23,42,0.9);
      }
      .muted { color: rgba(15,23,42,0.62); font-size: 12px; }

      /* Badges */
      .badges { display:flex; flex-wrap:wrap; gap:8px; margin-top: 8px; }
      .badge {
        font-size: 12px;
        padding: 7px 10px;
        border-radius: 999px;
        border: 1px solid rgba(15,23,42,0.10);
        background: rgba(15,23,42,0.03);
        color: rgba(15,23,42,0.70);
      }
      .badge b { color: rgba(15,23,42,0.92); }

      /* Streamlit widgets tweaks */
      div[data-testid="stFileUploader"] section {
        border-radius: 16px !important;
        border: 1px dashed rgba(15,23,42,0.22) !important;
        background: rgba(255,255,255,0.65) !important;
      }
      div[data-testid="stFileUploader"] section:hover {
        border-color: rgba(99,102,241,0.55) !important;
        box-shadow: 0 10px 20px rgba(99,102,241,0.10);
      }

      /* Buttons */
      .stButton button {
        border-radius: 14px !important;
        font-weight: 800 !important;
        padding: 10px 14px !important;
      }
      .stDownloadButton button {
        border-radius: 14px !important;
        font-weight: 800 !important;
        padding: 10px 14px !important;
      }

      /* Section spacing */
      .spacer { height: 10px; }
      hr { border: none; height: 1px; background: rgba(15,23,42,0.08); margin: 12px 0; }

    </style>
    """,
    unsafe_allow_html=True,
)

st.markdown(
    """
    <div class="hero">
      <div class="brandwrap">
        <div class="logo"></div>
        <div>
          <div class="brand">DocFlow Converter</div>
          <div class="tagline">Beautiful web converter • Upload → Choose → Convert → Download</div>
        </div>
      </div>
      <div class="chiprow">
        <div class="chip">100% Web (Streamlit)</div>
        <div class="chip">Cloud-safe OCR</div>
        <div class="chip">ZIP bundle exports</div>
      </div>
    </div>
    """,
    unsafe_allow_html=True,
)


# ============================================================
# HELPERS
# ============================================================
def now_stamp() -> str:
    return datetime.utcnow().strftime("%Y-%m-%d_%H%M%S_UTC")


def safe_filename(base: str) -> str:
    base = re.sub(r"[^A-Za-z0-9._-]+", "_", base).strip("_")
    return base[:120] if base else "output"


def build_zip(files: Dict[str, bytes]) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, mode="w", compression=zipfile.ZIP_DEFLATED) as z:
        for path, data in files.items():
            z.writestr(path, data)
    return buf.getvalue()


def infer_type(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        return "PDF"
    if ext in [".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"]:
        return "IMAGE"
    if ext in [".xlsx", ".xlsm"]:
        return "EXCEL"
    if ext == ".docx":
        return "WORD"
    if ext == ".pptx":
        return "PPT"
    return "UNKNOWN"


def mime_for(name: str) -> str:
    lname = name.lower()
    if lname.endswith(".zip"):
        return "application/zip"
    if lname.endswith(".pdf"):
        return "application/pdf"
    if lname.endswith(".docx"):
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if lname.endswith(".pptx"):
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if lname.endswith(".xlsx"):
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if lname.endswith(".csv"):
        return "text/csv"
    if lname.endswith(".json"):
        return "application/json"
    if lname.endswith(".txt"):
        return "text/plain"
    if lname.endswith(".png"):
        return "image/png"
    return "application/octet-stream"


# ============================================================
# OCR / PDF (cloud-safe: EasyOCR + PyMuPDF + pdfplumber)
# ============================================================
@st.cache_resource(show_spinner=False)
def _easyocr_reader(lang_code: str):
    import easyocr
    return easyocr.Reader([lang_code], gpu=False)


def _ui_lang_to_easyocr(ui_lang: str) -> str:
    return "en" if ui_lang == "eng" else "en"


def pdf_textlayer_extract(pdf_bytes: bytes, max_pages: int) -> List[str]:
    import pdfplumber
    texts = []
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for p in pdf.pages[:max_pages]:
            texts.append(p.extract_text() or "")
    return texts


def pdf_render_pages_to_images(pdf_bytes: bytes, dpi: int, max_pages: int) -> List[Image.Image]:
    import fitz
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    imgs: List[Image.Image] = []
    for i in range(min(max_pages, doc.page_count)):
        pix = doc.load_page(i).get_pixmap(dpi=int(dpi))
        imgs.append(Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB"))
    doc.close()
    return imgs


def ocr_image_to_text(pil_img: Image.Image, lang_ui: str = "eng") -> str:
    reader = _easyocr_reader(_ui_lang_to_easyocr(lang_ui))
    arr = np.array(pil_img.convert("RGB"))
    try:
        lines = reader.readtext(arr, detail=0, paragraph=True)
    except TypeError:
        lines = reader.readtext(arr, detail=0)
    return "\n".join([t.strip() for t in lines if t and str(t).strip()]).strip()


def pdf_hybrid_text_extract(pdf_bytes: bytes, max_pages: int, lang: str, dpi: int) -> List[str]:
    layer = pdf_textlayer_extract(pdf_bytes, max_pages=max_pages)
    needs_ocr = []
    for t in layer:
        t2 = re.sub(r"\s+", "", t or "")
        needs_ocr.append(len(t2) < 40)

    if not any(needs_ocr):
        return layer

    imgs = pdf_render_pages_to_images(pdf_bytes, dpi=dpi, max_pages=max_pages)
    out = []
    for i, base_text in enumerate(layer):
        if i < len(imgs) and needs_ocr[i]:
            try:
                t = ocr_image_to_text(imgs[i], lang_ui=lang)
                out.append(t if t else base_text)
            except Exception:
                out.append(base_text)
        else:
            out.append(base_text)
    return out


def pdf_to_images_zip(pdf_bytes: bytes, max_pages: int, dpi: int = 220) -> Tuple[bytes, int]:
    imgs = pdf_render_pages_to_images(pdf_bytes, dpi=dpi, max_pages=max_pages)
    files = {}
    for i, im in enumerate(imgs, start=1):
        buf = io.BytesIO()
        im.save(buf, format="PNG")
        files[f"pages/page_{i:03d}.png"] = buf.getvalue()
    return build_zip(files), len(imgs)


def pdf_metadata_to_json(pdf_bytes: bytes) -> bytes:
    from pypdf import PdfReader
    r = PdfReader(io.BytesIO(pdf_bytes))
    md = r.metadata or {}
    out = {"page_count": len(r.pages)}
    for k, v in md.items():
        out[str(k)] = str(v) if v is not None else None
    return json.dumps(out, ensure_ascii=False, indent=2).encode("utf-8")


# ============================================================
# TABLES (text-layer first, then OCR fallback)
# ============================================================
def normalize_cell_text_clean(val):
    if val is None:
        return val
    s = str(val).replace("\r\n", "\n").replace("\r", "\n")
    s = re.sub(r"\n+", " ", s)
    s = s.replace("\u00a0", " ")
    s = re.sub(r"[ \t]+", " ", s).strip()
    s = re.sub(r"^\|\s*", "", s)
    s = re.sub(r"\s*([,/:\.\-\+])\s*", r"\1", s)
    return s


def extract_tables_pdf_textlayer(pdf_bytes: bytes, max_pages: int) -> List[pd.DataFrame]:
    import pdfplumber
    dfs: List[pd.DataFrame] = []
    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages[:max_pages]:
            tbls = page.extract_tables()
            if not tbls:
                continue
            for t in tbls:
                if t:
                    df = pd.DataFrame(t)
                    df = df.replace("", np.nan).dropna(axis=0, how="all").dropna(axis=1, how="all").fillna("")
                    if not df.empty:
                        dfs.append(df)
    return dfs


def df_to_json_records(df: pd.DataFrame) -> List[Dict[str, Any]]:
    df2 = df.copy()
    cols = []
    for i, c in enumerate(df2.columns):
        name = str(c).strip() if c is not None else ""
        if not name or name.lower() in {"nan", "none"}:
            name = f"col_{i+1}"
        cols.append(name)
    df2.columns = cols
    df2 = df2.where(pd.notnull(df2), None)
    return df2.to_dict(orient="records")


def build_tables_bundle(tables: List[pd.DataFrame], base: str) -> Dict[str, bytes]:
    from openpyxl.styles import Alignment

    cleaned = [df.applymap(normalize_cell_text_clean) for df in tables]
    files: Dict[str, bytes] = {}

    excel_buf = io.BytesIO()
    with pd.ExcelWriter(excel_buf, engine="openpyxl") as writer:
        for i, df in enumerate(cleaned, start=1):
            df.to_excel(writer, sheet_name=f"Table_{i}"[:31], index=False, header=False)
        wb = writer.book
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    cell.alignment = Alignment(wrap_text=False, vertical="top")
    files[f"{base}.xlsx"] = excel_buf.getvalue()

    combined_json = {"tables": []}
    for i, df in enumerate(cleaned, start=1):
        files[f"{base}/tables/table_{i}.csv"] = df.to_csv(index=False, header=False).encode("utf-8")
        one = {"table_index": i, "rows": df_to_json_records(df)}
        files[f"{base}/tables/table_{i}.json"] = json.dumps(one, ensure_ascii=False, indent=2).encode("utf-8")
        combined_json["tables"].append(one)

    files[f"{base}/tables/combined.json"] = json.dumps(combined_json, ensure_ascii=False, indent=2).encode("utf-8")
    files[f"{base}/manifest.json"] = json.dumps({"type": "tables_export", "table_count": len(cleaned)}, indent=2).encode("utf-8")
    return files


def extract_table_from_image_webonly(
    img: Image.Image,
    lang_ui: str,
    min_conf_0_100: int,
    table_mode: str,
    enhance: bool,
    deskew: bool
) -> Tuple[List[pd.DataFrame], str]:
    """
    bordered: detect grid-ish cell boxes + OCR per cell (good for most invoices/line tables)
    borderless: OCR whole image -> split columns by multi-space (best-effort)
    """
    import cv2

    conf01 = max(0.10, min(0.95, float(min_conf_0_100) / 100.0))
    reader = _easyocr_reader(_ui_lang_to_easyocr(lang_ui))

    rgb = np.array(img.convert("RGB"))
    bgr = cv2.cvtColor(rgb, cv2.COLOR_RGB2BGR)
    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)

    if enhance:
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        gray = clahe.apply(gray)
        gray = cv2.GaussianBlur(gray, (3, 3), 0)

    bw = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY_INV, 25, 10)

    if deskew:
        coords = np.column_stack(np.where(bw > 0))
        if coords.size > 0:
            angle = cv2.minAreaRect(coords)[-1]
            if angle < -45:
                angle = -(90 + angle)
            else:
                angle = -angle
            h, w = bw.shape[:2]
            M = cv2.getRotationMatrix2D((w // 2, h // 2), angle, 1.0)
            bw = cv2.warpAffine(bw, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)

    if table_mode == "bordered":
        h, w = bw.shape[:2]
        hor_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (max(20, w // 30), 1))
        ver_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, max(20, h // 30)))

        horizontal = cv2.erode(bw, hor_kernel, iterations=1)
        horizontal = cv2.dilate(horizontal, hor_kernel, iterations=2)

        vertical = cv2.erode(bw, ver_kernel, iterations=1)
        vertical = cv2.dilate(vertical, ver_kernel, iterations=2)

        grid = cv2.addWeighted(horizontal, 0.5, vertical, 0.5, 0.0)
        grid = cv2.dilate(grid, cv2.getStructuringElement(cv2.MORPH_RECT, (3, 3)), iterations=1)

        contours, _ = cv2.findContours(grid, cv2.RETR_TREE, cv2.CHAIN_APPROX_SIMPLE)
        boxes = []
        for c in contours:
            x, y, ww, hh = cv2.boundingRect(c)
            if ww < 35 or hh < 18:
                continue
            if ww > w * 0.98 and hh > h * 0.98:
                continue
            area = ww * hh
            if area > (w * h) * 0.6:
                continue
            boxes.append((x, y, ww, hh))
        boxes = sorted(boxes, key=lambda b: (b[1], b[0]))

        if len(boxes) < 6:
            return [], "Bordered: not enough cells detected. Try Borderless or increase DPI/clarity."

        texts = []
        for (x, y, ww, hh) in boxes:
            crop = rgb[y:y+hh, x:x+ww]
            parts = reader.readtext(crop, detail=1)
            good = [p for p in parts if len(p) >= 3 and float(p[2]) >= conf01 and str(p[1]).strip()]
            if not good:
                texts.append("")
            else:
                good = sorted(good, key=lambda p: p[0][0][0])
                texts.append(" ".join([str(p[1]).strip() for p in good]).strip())

        centers = np.array([(x + ww / 2, y + hh / 2) for (x, y, ww, hh) in boxes], dtype=float)
        xs = centers[:, 0].tolist()
        ys = centers[:, 1].tolist()

        def cluster(vals: List[float], gap: float) -> List[int]:
            order = np.argsort(vals)
            cid = np.zeros(len(vals), dtype=int)
            cur = 0
            prev = None
            for idx in order:
                v = vals[idx]
                if prev is None:
                    cid[idx] = cur
                    prev = v
                    continue
                if abs(v - prev) > gap:
                    cur += 1
                cid[idx] = cur
                prev = v
            return cid.tolist()

        row_ids = cluster(ys, gap=18.0)
        col_ids = cluster(xs, gap=28.0)

        n_rows = max(row_ids) + 1
        n_cols = max(col_ids) + 1
        grid_cells = [["" for _ in range(n_cols)] for _ in range(n_rows)]
        for i, t in enumerate(texts):
            r = row_ids[i]
            c = col_ids[i]
            grid_cells[r][c] = (grid_cells[r][c] + " " + t).strip() if grid_cells[r][c] else t

        df = pd.DataFrame(grid_cells)
        df = df.replace("", np.nan).dropna(axis=0, how="all").dropna(axis=1, how="all").fillna("")
        if df.empty:
            return [], "Bordered: detected cells but output was empty after cleaning."
        return [df], f"Bordered: extracted table {df.shape[0]}×{df.shape[1]}."

    # Borderless mode
    parts = reader.readtext(rgb, detail=1)
    good = [p for p in parts if len(p) >= 3 and float(p[2]) >= conf01 and str(p[1]).strip()]
    if not good:
        return [], "Borderless: no OCR text found above confidence threshold."

    good = sorted(good, key=lambda p: (p[0][0][1], p[0][0][0]))
    lines: List[str] = []
    cur = []
    last_y = None
    for box, txt, conf in good:
        y = float(box[0][1])
        if last_y is None or abs(y - last_y) < 14:
            cur.append(str(txt).strip())
        else:
            lines.append(" ".join(cur).strip())
            cur = [str(txt).strip()]
        last_y = y
    if cur:
        lines.append(" ".join(cur).strip())

    rows = [re.split(r"\s{2,}", ln.strip()) for ln in lines if ln.strip()]
    max_cols = max((len(r) for r in rows), default=0)
    padded = [r + [""] * (max_cols - len(r)) for r in rows]
    df = pd.DataFrame(padded).replace("", np.nan).dropna(axis=0, how="all").dropna(axis=1, how="all").fillna("")
    if df.empty:
        return [], "Borderless: OCR found text but failed to build a table."
    return [df], f"Borderless: extracted table {df.shape[0]}×{df.shape[1]} (best-effort)."


# ============================================================
# SEARCHABLE PDF (web-only overlay text)
# ============================================================
def make_searchable_pdf_from_images(images: List[Image.Image], lang_ui: str, conf01: float) -> bytes:
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.colors import Color

    reader = _easyocr_reader(_ui_lang_to_easyocr(lang_ui))
    invisible = Color(0, 0, 0, alpha=0.01)

    buf = io.BytesIO()
    c = canvas.Canvas(buf)

    for img in images:
        w_px, h_px = img.size
        c.setPageSize((w_px, h_px))
        c.drawImage(ImageReader(img), 0, 0, width=w_px, height=h_px, mask="auto")

        arr = np.array(img.convert("RGB"))
        results = reader.readtext(arr, detail=1)

        c.setFillColor(invisible)
        for box, text, conf in results:
            if float(conf) < conf01:
                continue
            if not text or not str(text).strip():
                continue

            xs = [p[0] for p in box]
            ys = [p[1] for p in box]
            x_min, x_max = float(min(xs)), float(max(xs))
            y_min, y_max = float(min(ys)), float(max(ys))

            pdf_x = x_min
            pdf_y = h_px - y_max
            font_size = max(6.0, min(24.0, (y_max - y_min) * 0.8))
            c.setFont("Helvetica", font_size)
            c.drawString(pdf_x, pdf_y, re.sub(r"\s+", " ", str(text))[:200])

        c.showPage()

    c.save()
    return buf.getvalue()


# ============================================================
# Office conversions
# ============================================================
def excel_to_pdf_bytes(xlsx_bytes: bytes, max_rows: int = 90, max_cols: int = 14) -> bytes:
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet

    styles = getSampleStyleSheet()
    story = []

    xls = pd.ExcelFile(io.BytesIO(xlsx_bytes))
    for si, sheet in enumerate(xls.sheet_names, start=1):
        df = xls.parse(sheet).iloc[:max_rows, :max_cols]
        story.append(Paragraph(f"{sheet}", styles["Heading2"]))
        story.append(Spacer(1, 8))
        data = [list(df.columns.astype(str))] + df.astype(str).values.tolist()
        tbl = Table(data, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f2937")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
        ]))
        story.append(tbl)
        if si < len(xls.sheet_names):
            story.append(PageBreak())

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(A4), leftMargin=18, rightMargin=18, topMargin=18, bottomMargin=18)
    doc.build(story)
    return buf.getvalue()


def excel_to_word_docx_bytes(xlsx_bytes: bytes, max_rows: int = 120, max_cols: int = 14) -> bytes:
    from docx import Document
    from docx.shared import Pt

    xls = pd.ExcelFile(io.BytesIO(xlsx_bytes))
    doc = Document()
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    for si, sheet in enumerate(xls.sheet_names, start=1):
        df = xls.parse(sheet).iloc[:max_rows, :max_cols]
        doc.add_heading(sheet, level=2)

        table = doc.add_table(rows=1, cols=len(df.columns))
        hdr = table.rows[0].cells
        for j, c in enumerate(df.columns.astype(str).tolist()):
            hdr[j].text = c

        for _, row in df.astype(str).iterrows():
            cells = table.add_row().cells
            for j, val in enumerate(row.tolist()):
                cells[j].text = val

        if si < len(xls.sheet_names):
            doc.add_paragraph("")

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def word_to_excel_tables(docx_bytes: bytes) -> Optional[bytes]:
    from docx import Document
    doc = Document(io.BytesIO(docx_bytes))
    if not doc.tables:
        return None
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="openpyxl") as writer:
        for i, t in enumerate(doc.tables, start=1):
            rows = []
            for r in t.rows:
                rows.append([c.text.strip() for c in r.cells])
            pd.DataFrame(rows).to_excel(writer, sheet_name=f"Table_{i}"[:31], index=False, header=False)
    return out.getvalue()


def docx_to_plain_text(docx_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(docx_bytes))
    parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(parts).strip()


def pdf_to_docx_high_fidelity(pdf_bytes: bytes) -> bytes:
    from pdf2docx import Converter
    tmp_id = datetime.utcnow().strftime("%Y%m%d_%H%M%S_%f")
    in_path = f"/tmp/in_{tmp_id}.pdf"
    out_path = f"/tmp/out_{tmp_id}.docx"
    with open(in_path, "wb") as f:
        f.write(pdf_bytes)
    try:
        cv = Converter(in_path)
        cv.convert(out_path, start=0, end=None)
        cv.close()
        with open(out_path, "rb") as f:
            return f.read()
    finally:
        for p in (in_path, out_path):
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass


# ============================================================
# STATE
# ============================================================
if "outputs" not in st.session_state:
    st.session_state.outputs = {}  # name -> bytes

if "history" not in st.session_state:
    st.session_state.history = []  # list[dict]

if "last_preview_key" not in st.session_state:
    st.session_state.last_preview_key = None


def push_history(task_label: str, task_key: str, fname: str):
    st.session_state.history.insert(0, {
        "time": datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S UTC"),
        "task": task_label,
        "task_key": task_key,
        "file": fname,
    })
    st.session_state.history = st.session_state.history[:60]


# ============================================================
# SIDEBAR (kept minimal, premium)
# ============================================================
with st.sidebar:
    st.markdown("### Settings")
    with st.expander("Advanced options", expanded=False):
        ocr_lang = st.selectbox("OCR language", ["eng"], index=0)
        max_pages = st.slider("Max pages (PDF)", 1, 80, 12)
        ocr_dpi = st.slider("OCR DPI", 200, 400, 260, step=10)
        min_conf = st.slider("OCR confidence", 10, 95, 50)
        table_mode = st.selectbox("Table mode", ["bordered", "borderless"], index=0)
        enhance = st.checkbox("Enhance scan", value=True)
        deskew = st.checkbox("Deskew scan", value=True)

    st.markdown("---")
    st.caption("Tip: **Bordered** works best for invoices & grid tables. Use **Borderless** for no-lines tables.")


# ============================================================
# CONVERSION MENU
# ============================================================
TASKS_BY_TYPE = {
    "PDF": [
        ("Extract Tables → Excel/CSV/JSON (ZIP)", "pdf_tables"),
        ("Extract Text (Hybrid) → TXT", "pdf_text_txt"),
        ("PDF → Editable Word (DOCX)", "pdf_to_docx"),
        ("Create Searchable PDF (OCR layer)", "pdf_searchable"),
        ("Pages → PNG (ZIP)", "pdf_pages_png"),
        ("Metadata → JSON", "pdf_meta_json"),
    ],
    "IMAGE": [
        ("OCR Image → TXT", "img_text_txt"),
        ("Image Table → Excel/CSV/JSON (ZIP)", "img_tables"),
        ("Image → PDF", "img_to_pdf"),
        ("Image → Searchable PDF (OCR layer)", "img_searchable_pdf"),
    ],
    "EXCEL": [
        ("Excel → PDF", "xlsx_to_pdf"),
        ("Excel → Word (DOCX)", "xlsx_to_docx"),
    ],
    "WORD": [
        ("Word → TXT", "docx_to_txt"),
        ("Word Tables → Excel (XLSX)", "docx_tables_to_xlsx"),
    ],
    "PPT": [
        ("PPT Text → TXT/JSON (ZIP)", "pptx_text_bundle"),
        ("PPT Images → ZIP", "pptx_images_zip"),
    ],
}


# ============================================================
# MAIN LAYOUT (polished)
# ============================================================
left, right = st.columns([1.25, 1.0], gap="large")

with left:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.markdown('<div class="cardtitle"><h3>Upload & Choose</h3><span class="muted">Step 1–3</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="muted">Drop a file below. You’ll see the best available conversions instantly.</div>', unsafe_allow_html=True)
    st.markdown('<div class="spacer"></div>', unsafe_allow_html=True)

    uploaded = st.file_uploader(
        "Upload file",
        type=["pdf", "png", "jpg", "jpeg", "webp", "tif", "tiff", "bmp", "docx", "xlsx", "xlsm", "pptx"],
        label_visibility="collapsed",
    )

    st.markdown("<hr/>", unsafe_allow_html=True)

    if uploaded:
        filename = uploaded.name
        file_bytes = uploaded.read()
        ftype = infer_type(filename)
        base = safe_filename(os.path.splitext(filename)[0])

        st.markdown(
            f"""
            <div class="badges">
              <div class="badge"><b>Type</b>: {ftype}</div>
              <div class="badge"><b>Name</b>: {filename}</div>
              <div class="badge"><b>Size</b>: {len(file_bytes)/1024:.1f} KB</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

        # Show beautiful static "capability" badges
        caps = {
            "PDF": ["PDF → Word", "PDF → Tables", "PDF → Searchable", "PDF → Images", "PDF → Metadata"],
            "IMAGE": ["Image → OCR", "Image → Tables", "Image → PDF", "Image → Searchable"],
            "EXCEL": ["Excel → PDF", "Excel → Word"],
            "WORD": ["Word → Text", "Word → Tables"],
            "PPT": ["PPT → Text", "PPT → Images"],
        }.get(ftype, ["Upload a supported file"])

        st.markdown(
            "<div class='badges'>" +
            "".join([f"<div class='badge'>{c}</div>" for c in caps]) +
            "</div>",
            unsafe_allow_html=True
        )

        st.markdown("<div class='spacer'></div>", unsafe_allow_html=True)

        opts = TASKS_BY_TYPE.get(ftype, [])
        task_labels = [t[0] for t in opts] if opts else ["No conversions available"]
        task_disabled = not bool(opts)

        task_label = st.selectbox("Conversion", task_labels, index=0, disabled=task_disabled)
        task_key = dict(opts).get(task_label) if not task_disabled else None

    else:
        filename = None
        file_bytes = None
        ftype = "—"
        base = "output"
        task_label = None
        task_key = None

        st.markdown(
            """
            <div class="badges">
              <div class="badge"><b>Supported</b>: PDF, Images, Word, Excel, PPT</div>
              <div class="badge">OCR Tables → Excel</div>
              <div class="badge">PDF → Editable Word</div>
              <div class="badge">Batch ZIP Downloads</div>
            </div>
            """,
            unsafe_allow_html=True,
        )

    st.markdown("</div>", unsafe_allow_html=True)


with right:
    st.markdown('<div class="card">', unsafe_allow_html=True)
    st.markdown('<div class="cardtitle"><h3>Convert & Download</h3><span class="muted">Step 4–5</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="muted">Your outputs appear here with preview + one-click ZIP bundle.</div>', unsafe_allow_html=True)
    st.markdown('<div class="spacer"></div>', unsafe_allow_html=True)

    convert_disabled = not (uploaded and task_key)
    colA, colB = st.columns([0.70, 0.30])
    with colA:
        convert_btn = st.button("Convert", type="primary", disabled=convert_disabled, use_container_width=True)
    with colB:
        clear_btn = st.button("Clear outputs", disabled=not bool(st.session_state.outputs), use_container_width=True)
        if clear_btn:
            st.session_state.outputs = {}
            st.session_state.last_preview_key = None
            st.rerun()

    st.markdown("<hr/>", unsafe_allow_html=True)

    # Outputs + Preview tabs
    tab_files, tab_preview = st.tabs(["📦 Files", "👁️ Preview"])

    with tab_files:
        if not st.session_state.outputs:
            st.info("No outputs yet. Upload a file, choose a conversion, and click **Convert**.")
        else:
            st.success("Conversion complete. Download below.")

            # individual downloads
            for out_name, out_bytes in st.session_state.outputs.items():
                st.download_button(
                    label=f"Download {out_name}",
                    data=out_bytes,
                    file_name=out_name,
                    mime=mime_for(out_name),
                    use_container_width=True,
                    key=f"dl_{out_name}_{now_stamp()}",
                )

            # bundle zip
            zname = f"{safe_filename(os.path.splitext(filename)[0])}_bundle_{now_stamp()}.zip" if filename else f"bundle_{now_stamp()}.zip"
            st.download_button(
                label="Download ALL as ZIP",
                data=build_zip(st.session_state.outputs),
                file_name=zname,
                mime="application/zip",
                use_container_width=True,
                key=f"dl_zip_{now_stamp()}",
            )

    with tab_preview:
        if not st.session_state.outputs:
            st.caption("Nothing to preview yet.")
        else:
            # choose preview file
            out_names = list(st.session_state.outputs.keys())
            default_idx = 0
            if st.session_state.last_preview_key in out_names:
                default_idx = out_names.index(st.session_state.last_preview_key)

            pick = st.selectbox("Preview file", out_names, index=default_idx)
            st.session_state.last_preview_key = pick
            data = st.session_state.outputs[pick]

            low = pick.lower()
            if low.endswith(".txt"):
                try:
                    st.text_area("Text", data.decode("utf-8", errors="replace"), height=320)
                except Exception:
                    st.caption("Cannot decode this TXT file.")
            elif low.endswith(".json"):
                try:
                    st.json(json.loads(data.decode("utf-8", errors="replace")))
                except Exception:
                    st.caption("Cannot parse JSON.")
            elif low.endswith(".png"):
                st.image(Image.open(io.BytesIO(data)))
            elif low.endswith(".pdf"):
                st.caption("PDF generated. (Streamlit preview varies by browser; download to view.)")
            elif low.endswith(".xlsx"):
                st.caption("Excel generated. Showing a small preview of first sheet (best-effort).")
                try:
                    xls = pd.ExcelFile(io.BytesIO(data))
                    df = xls.parse(xls.sheet_names[0]).head(25)
                    st.dataframe(df, use_container_width=True)
                except Exception:
                    st.caption("Preview not available for this XLSX.")
            else:
                st.caption("Preview not supported for this file type. Download it.")

    st.markdown("</div>", unsafe_allow_html=True)


# ============================================================
# CONVERSION EXECUTION
# ============================================================
if convert_btn and uploaded and task_key and file_bytes and filename:
    st.session_state.outputs = {}
    base = safe_filename(os.path.splitext(filename)[0])

    try:
        outputs: Dict[str, bytes] = {}

        # PDF
        if task_key == "pdf_text_txt":
            pages = pdf_hybrid_text_extract(file_bytes, max_pages=max_pages, lang=ocr_lang, dpi=ocr_dpi)
            txt = "\n\n".join([p.strip() for p in pages if p is not None]).strip()
            outputs[f"{base}.txt"] = (txt + "\n").encode("utf-8")

        elif task_key == "pdf_to_docx":
            # High fidelity first; fallback to hybrid text -> docx paragraphs
            try:
                docx_bytes = pdf_to_docx_high_fidelity(file_bytes)
                outputs[f"{base}.docx"] = docx_bytes
            except Exception:
                from docx import Document
                pages = pdf_hybrid_text_extract(file_bytes, max_pages=max_pages, lang=ocr_lang, dpi=ocr_dpi)
                doc = Document()
                for p in pages:
                    p = (p or "").strip()
                    if p:
                        for line in p.splitlines():
                            if line.strip():
                                doc.add_paragraph(line.strip())
                        doc.add_paragraph("")
                out = io.BytesIO()
                doc.save(out)
                outputs[f"{base}.docx"] = out.getvalue()

        elif task_key == "pdf_tables":
            tables = []
            try:
                tables = extract_tables_pdf_textlayer(file_bytes, max_pages=max_pages)
            except Exception:
                tables = []

            # OCR fallback for scanned PDFs
            if not tables:
                imgs = pdf_render_pages_to_images(file_bytes, dpi=ocr_dpi, max_pages=max_pages)
                ocr_tables: List[pd.DataFrame] = []
                for im in imgs:
                    tbs, _lg = extract_table_from_image_webonly(
                        im, lang_ui=ocr_lang, min_conf_0_100=min_conf,
                        table_mode=table_mode, enhance=enhance, deskew=deskew
                    )
                    ocr_tables.extend(tbs)
                tables = ocr_tables

            if not tables:
                raise RuntimeError("No tables found. Try increasing DPI, switching Table mode, or use a clearer scan.")

            root = f"{base}_tables_{now_stamp()}"
            bundle = build_tables_bundle(tables, base=root)
            outputs[f"{root}.xlsx"] = bundle[f"{root}.xlsx"]
            outputs[f"{root}.zip"] = build_zip(bundle)

        elif task_key == "pdf_pages_png":
            z, _ = pdf_to_images_zip(file_bytes, max_pages=max_pages, dpi=220)
            outputs[f"{base}_pages_{now_stamp()}.zip"] = z

        elif task_key == "pdf_meta_json":
            outputs[f"{base}_metadata.json"] = pdf_metadata_to_json(file_bytes)

        elif task_key == "pdf_searchable":
            conf01 = max(0.10, min(0.95, float(min_conf) / 100.0))
            imgs = pdf_render_pages_to_images(file_bytes, dpi=ocr_dpi, max_pages=max_pages)
            outputs[f"{base}_searchable.pdf"] = make_searchable_pdf_from_images(imgs, lang_ui=ocr_lang, conf01=conf01)

        # IMAGE
        elif task_key == "img_text_txt":
            im = Image.open(io.BytesIO(file_bytes)).convert("RGB")
            txt = ocr_image_to_text(im, lang_ui=ocr_lang)
            outputs[f"{base}.txt"] = (txt + "\n").encode("utf-8")

        elif task_key == "img_to_pdf":
            im = Image.open(io.BytesIO(file_bytes)).convert("RGB")
            buf = io.BytesIO()
            im.save(buf, format="PDF")
            outputs[f"{base}.pdf"] = buf.getvalue()

        elif task_key == "img_searchable_pdf":
            conf01 = max(0.10, min(0.95, float(min_conf) / 100.0))
            im = Image.open(io.BytesIO(file_bytes)).convert("RGB")
            outputs[f"{base}_searchable.pdf"] = make_searchable_pdf_from_images([im], lang_ui=ocr_lang, conf01=conf01)

        elif task_key == "img_tables":
            im = Image.open(io.BytesIO(file_bytes)).convert("RGB")
            tables, _lg = extract_table_from_image_webonly(
                im, lang_ui=ocr_lang, min_conf_0_100=min_conf,
                table_mode=table_mode, enhance=enhance, deskew=deskew
            )
            if not tables:
                raise RuntimeError("No tables found. Try Borderless mode, enable Enhance/Deskew, or use a clearer image.")

            root = f"{base}_tables_{now_stamp()}"
            bundle = build_tables_bundle(tables, base=root)
            outputs[f"{root}.xlsx"] = bundle[f"{root}.xlsx"]
            outputs[f"{root}.zip"] = build_zip(bundle)

        # EXCEL
        elif task_key == "xlsx_to_pdf":
            outputs[f"{base}.pdf"] = excel_to_pdf_bytes(file_bytes)

        elif task_key == "xlsx_to_docx":
            outputs[f"{base}.docx"] = excel_to_word_docx_bytes(file_bytes)

        # WORD
        elif task_key == "docx_to_txt":
            txt = docx_to_plain_text(file_bytes)
            outputs[f"{base}.txt"] = (txt + "\n").encode("utf-8")

        elif task_key == "docx_tables_to_xlsx":
            xlsx = word_to_excel_tables(file_bytes)
            if xlsx is None:
                raise RuntimeError("No tables found in this Word document.")
            outputs[f"{base}_tables.xlsx"] = xlsx

        # PPT
        elif task_key == "pptx_text_bundle":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides = []
            all_text = []
            for i, slide in enumerate(prs.slides, start=1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = shape.text.strip()
                        if t:
                            parts.append(t)
                stext = "\n".join(parts).strip()
                slides.append({"slide": i, "text": stext})
                all_text.append(f"Slide {i}\n{stext}".strip())

            files = {
                "slides.txt": ("\n\n".join(all_text).strip() + "\n").encode("utf-8"),
                "slides.json": json.dumps({"slides": slides}, ensure_ascii=False, indent=2).encode("utf-8"),
            }
            outputs[f"{base}_ppt_export_{now_stamp()}.zip"] = build_zip(files)

        elif task_key == "pptx_images_zip":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            files = {}
            count = 0
            for si, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == 13:
                            image = shape.image
                            ext = (image.ext or "bin").lower()
                            count += 1
                            files[f"images/slide_{si:03d}_{count:03d}.{ext}"] = image.blob
                    except Exception:
                        pass
            if not files:
                raise RuntimeError("No embedded images found in this PPTX.")
            outputs[f"{base}_images_{now_stamp()}.zip"] = build_zip(files)

        else:
            raise RuntimeError("Conversion not implemented.")

        st.session_state.outputs = outputs
        push_history(task_label=task_label, task_key=task_key, fname=filename)

        st.success("Done — outputs are ready in the Files tab.")

    except Exception as e:
        st.session_state.outputs = {}
        st.error(str(e))


# ============================================================
# RECENT CONVERSIONS (clean + optional)
# ============================================================
st.markdown("<div class='spacer'></div>", unsafe_allow_html=True)

with st.expander("Recent conversions (optional)", expanded=False):
    if not st.session_state.history:
        st.caption("No history yet.")
    else:
        dfh = pd.DataFrame(st.session_state.history)
        dfh = dfh[["time", "file", "task"]]
        st.dataframe(dfh, use_container_width=True, hide_index=True)
        st.caption("History is informational. It does not auto-run anything.")


